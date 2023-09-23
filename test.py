from pptx import Presentation

path = "C:/Users/georg/OneDrive/Documents/Rice University/Projects/HackRice 2023/PrezPal/.slides/Whale Tank Details.pptx"
path = ".slides/Whale Tank Details.pptx"
prs = Presentation(path)
for slide in prs.slides:

    print([val for val in dir(slide) if val[0] != "_"])

    for shape in slide.shapes:
        print([val for val in dir(shape) if val[0] != "_"])

        for att in ["placeholder_format", "text", "image", "top", "left", "height", "width"]:
            if hasattr(shape, att):
                print(dir(getattr(shape, "placeholder_format")))
                print(getattr(shape, att))